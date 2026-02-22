"""
PDF Indexer — extracts text from PDFs page-by-page, generates embeddings,
and stores them in ChromaDB for semantic search.

Designed to be modular: add new file processors (e.g. PPTX) alongside
the PDF processor without refactoring.
"""

import hashlib
import os
from pathlib import Path

import fitz  # PyMuPDF
import chromadb
import ollama

EMBED_MODEL = "nomic-embed-text"
CHROMA_DIR = os.path.join(os.path.dirname(__file__), "chroma_data")
COLLECTION_NAME = "slides"
SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt"}


def extract_pages_from_pdf(pdf_path: str) -> list[dict]:
    """Extract text from each page of a PDF file.

    Returns a list of dicts: {page_number (1-indexed), text}
    """
    pages = []
    doc = fitz.open(pdf_path)
    for i in range(len(doc)):
        text = doc[i].get_text().strip()
        if text:  # skip empty pages
            pages.append({
                "page_number": i + 1,  # 1-indexed for display
                "text": text,
            })
    doc.close()
    return pages


def extract_pages_from_pptx(pptx_path: str) -> list[dict]:
    """Extract text from each slide of a PPTX file."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        text = "\n".join(texts).strip()
        if text:
            pages.append({"page_number": i, "text": text})
    return pages


def extract_pages_from_docx(docx_path: str) -> list[dict]:
    """Extract text from a DOCX file as a single 'page'."""
    from docx import Document
    doc = Document(docx_path)
    text = "\n".join(p.text for p in doc.paragraphs).strip()
    if text:
        return [{"page_number": 1, "text": text}]
    return []


def extract_pages_from_txt(txt_path: str) -> list[dict]:
    """Read a plain text file as a single 'page'."""
    text = Path(txt_path).read_text(encoding="utf-8", errors="replace").strip()
    if text:
        return [{"page_number": 1, "text": text}]
    return []


EXTRACTORS = {
    ".pdf": extract_pages_from_pdf,
    ".pptx": extract_pages_from_pptx,
    ".docx": extract_pages_from_docx,
    ".txt": extract_pages_from_txt,
}


def scan_folder(folder_path: str) -> list[dict]:
    """Recursively scan a folder for supported files and extract pages.

    Returns a list of dicts: {file_path (relative), page_number, text}
    """
    root = Path(folder_path)
    all_pages = []

    for file_path in sorted(root.rglob("*")):
        if file_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        if not file_path.is_file():
            continue

        relative_path = str(file_path.relative_to(root))
        extractor = EXTRACTORS.get(file_path.suffix.lower())
        if not extractor:
            continue
        pages = extractor(str(file_path))

        for page in pages:
            all_pages.append({
                "file_path": relative_path,
                "page_number": page["page_number"],
                "text": page["text"],
            })

    return all_pages


def generate_embeddings(texts: list[str]) -> list[list[float]]:
    """Generate embeddings for a list of texts using Ollama."""
    result = ollama.embed(model=EMBED_MODEL, input=texts)
    return result.embeddings


def _folder_fingerprint(folder_path: str) -> str:
    """Compute a hash of file paths + modification times in a folder."""
    root = Path(folder_path)
    entries = []
    for fp in sorted(root.rglob("*")):
        if fp.suffix.lower() in SUPPORTED_EXTENSIONS and fp.is_file():
            entries.append(f"{fp}:{fp.stat().st_mtime_ns}")
    return hashlib.sha256("\n".join(entries).encode()).hexdigest()[:16]


def index_folder(folder_path: str) -> dict:
    """Index all supported files in a folder into ChromaDB.

    Returns a summary dict with counts. Skips re-indexing if unchanged.
    """
    fingerprint = _folder_fingerprint(folder_path)

    # Check if already indexed with same fingerprint
    client = chromadb.PersistentClient(path=CHROMA_DIR)
    try:
        collection = client.get_collection(COLLECTION_NAME)
        root_doc = collection.get(ids=["__root__"])
        if root_doc["metadatas"]:
            stored_fp = root_doc["metadatas"][0].get("fingerprint")
            stored_root = root_doc["metadatas"][0].get("root_folder")
            if stored_fp == fingerprint and stored_root == folder_path:
                # Count existing documents (minus the __root__ marker)
                total = collection.count() - 1
                files = set()
                # Get all unique file paths from metadata
                all_meta = collection.get(where={"file_path": {"$ne": "__root__"}})
                for m in all_meta["metadatas"]:
                    files.add(m["file_path"])
                return {
                    "status": "ok",
                    "total_pages": total,
                    "total_files": len(files),
                    "files": sorted(files),
                    "message": "Already indexed (no changes detected)",
                }
    except Exception:
        pass

    pages = scan_folder(folder_path)
    if not pages:
        return {"status": "error", "message": "No supported files found", "total_pages": 0}

    # Delete existing collection if re-indexing
    try:
        client.delete_collection(COLLECTION_NAME)
    except Exception:
        pass

    collection = client.create_collection(
        name=COLLECTION_NAME,
        metadata={"hnsw:space": "cosine"},
    )

    # Store the root folder path so we can resolve files later
    # (stored as collection metadata isn't supported for arbitrary keys,
    #  so we store it in a special document)
    collection.add(
        ids=["__root__"],
        documents=["__root_folder__"],
        metadatas=[{"file_path": "__root__", "page_number": 0, "root_folder": folder_path, "fingerprint": fingerprint}],
        embeddings=[generate_embeddings(["root folder marker"])[0]],
    )

    # Batch embed and insert (batch size to avoid overwhelming Ollama)
    batch_size = 32
    total = len(pages)
    files_seen = set()

    for start in range(0, total, batch_size):
        batch = pages[start : start + batch_size]
        texts = [p["text"] for p in batch]
        embeddings = generate_embeddings(texts)

        ids = []
        documents = []
        metadatas = []

        for i, page in enumerate(batch):
            doc_id = f"{page['file_path']}::page_{page['page_number']}"
            ids.append(doc_id)
            documents.append(page["text"])
            metadatas.append({
                "file_path": page["file_path"],
                "page_number": page["page_number"],
            })
            files_seen.add(page["file_path"])

        collection.add(
            ids=ids,
            documents=documents,
            metadatas=metadatas,
            embeddings=embeddings,
        )

    return {
        "status": "ok",
        "total_pages": total,
        "total_files": len(files_seen),
        "files": sorted(files_seen),
    }


if __name__ == "__main__":
    # Quick test with the 170_Theory folder
    import json

    folder = "/Users/porter/Documents/170_Theory"
    print(f"Indexing: {folder}")
    result = index_folder(folder)
    print(json.dumps(result, indent=2))
